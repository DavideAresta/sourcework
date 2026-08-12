"""Turn document bytes into locator-annotated plain text.

Every extractor returns a list of ``(locator, text)`` blocks rather than one
big string. The locator is what ends up in the PRD's traceability table, so a
reader can open page 12 of the RFP and check the claim themselves.

Parser libraries are imported lazily and degrade to a clear warning if the
optional ``ingest`` extra is not installed.
"""

from __future__ import annotations

import csv
import io
import logging
import re
from collections.abc import Iterator

logger = logging.getLogger(__name__)

Block = tuple[str, str]  # (locator, text)


class UnsupportedDocument(RuntimeError):
    pass


def extract(data: bytes, media_type: str, filename: str = "") -> tuple[list[Block], list[str]]:
    """Dispatch on media type / extension. Returns ``(blocks, warnings)``."""
    warnings: list[str] = []
    kind = _classify(media_type, filename)
    try:
        extractor = _EXTRACTORS[kind]
    except KeyError:
        raise UnsupportedDocument(
            f"No extractor for {media_type!r} ({filename!r}). "
            "Supported: pdf, docx, pptx, xlsx, csv, html, markdown, text."
        ) from None

    try:
        blocks = list(extractor(data))
    except ImportError as exc:
        raise UnsupportedDocument(
            f"Parsing {kind} needs an optional dependency: {exc}. "
            'Install with: pip install "prdforge[ingest]"'
        ) from exc

    blocks = [(loc, _tidy(text)) for loc, text in blocks if text and text.strip()]
    if not blocks:
        warnings.append(f"{filename or kind}: no extractable text (scanned image? OCR not wired up)")
    return blocks, warnings


# ---------------------------------------------------------------------------


def _classify(media_type: str, filename: str) -> str:
    mt = (media_type or "").lower()
    name = filename.lower()
    if "pdf" in mt or name.endswith(".pdf"):
        return "pdf"
    if "wordprocessingml" in mt or name.endswith(".docx"):
        return "docx"
    if "presentationml" in mt or name.endswith(".pptx"):
        return "pptx"
    if "spreadsheetml" in mt or name.endswith((".xlsx", ".xlsm")):
        return "xlsx"
    if "csv" in mt or name.endswith((".csv", ".tsv")):
        return "csv"
    if "html" in mt or name.endswith((".html", ".htm")):
        return "html"
    if name.endswith((".md", ".markdown")) or "markdown" in mt:
        return "markdown"
    return "text"


def _tidy(text: str) -> str:
    text = text.replace(" ", " ")
    text = re.sub(r"[ \t]+", " ", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _pdf(data: bytes) -> Iterator[Block]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    for i, page in enumerate(reader.pages, start=1):
        yield f"p.{i}", page.extract_text() or ""


def _docx(data: bytes) -> Iterator[Block]:
    import docx

    doc = docx.Document(io.BytesIO(data))
    heading = "body"
    buffer: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "").lower()
        if style.startswith("heading") and para.text.strip():
            if buffer:
                yield heading, "\n".join(buffer)
                buffer = []
            heading = para.text.strip()
            yield f"heading: {heading}", para.text.strip()
        elif para.text.strip():
            buffer.append(para.text.strip())
    if buffer:
        yield heading, "\n".join(buffer)

    for t_index, table in enumerate(doc.tables, start=1):
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in table.rows]
        yield f"table {t_index}", "\n".join(rows)


def _pptx(data: bytes) -> Iterator[Block]:
    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    for i, slide in enumerate(deck.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            chunks.append(f"[speaker notes] {slide.notes_slide.notes_text_frame.text.strip()}")
        yield f"slide {i}", "\n".join(chunks)


def _xlsx(data: bytes) -> Iterator[Block]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
            if len(rows) >= 500:
                rows.append("... (truncated at 500 rows)")
                break
        yield f"sheet {sheet.title}", "\n".join(rows)


def _csv(data: bytes) -> Iterator[Block]:
    text = data.decode("utf-8", errors="replace")
    dialect = csv.Sniffer().sniff(text[:4096]) if text[:4096].strip() else csv.excel
    reader = csv.reader(io.StringIO(text), dialect)
    rows = [" | ".join(r) for r in reader]
    for start in range(0, len(rows), 200):
        yield f"rows {start + 1}-{min(start + 200, len(rows))}", "\n".join(rows[start : start + 200])


def _html(data: bytes) -> Iterator[Block]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data.decode("utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    current = "body"
    buffer: list[str] = []
    for el in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "td", "th", "pre"]):
        text = el.get_text(" ", strip=True)
        if not text:
            continue
        if el.name.startswith("h"):
            if buffer:
                yield current, "\n".join(buffer)
                buffer = []
            current = text
        else:
            buffer.append(text)
    if buffer:
        yield current, "\n".join(buffer)


def _markdown(data: bytes) -> Iterator[Block]:
    text = data.decode("utf-8", errors="replace")
    current = "preamble"
    buffer: list[str] = []
    for line in text.splitlines():
        if line.startswith("#"):
            if buffer:
                yield current, "\n".join(buffer)
                buffer = []
            current = line.lstrip("# ").strip() or current
        else:
            buffer.append(line)
    if buffer:
        yield current, "\n".join(buffer)


def _text(data: bytes) -> Iterator[Block]:
    text = data.decode("utf-8", errors="replace")
    paragraphs = [p for p in re.split(r"\n\s*\n", text) if p.strip()]
    if not paragraphs:
        return
    for i, para in enumerate(paragraphs, start=1):
        yield f"para {i}", para


_EXTRACTORS = {
    "pdf": _pdf,
    "docx": _docx,
    "pptx": _pptx,
    "xlsx": _xlsx,
    "csv": _csv,
    "html": _html,
    "markdown": _markdown,
    "text": _text,
}


def chunk(blocks: list[Block], max_chars: int = 12000) -> list[list[Block]]:
    """Group blocks into LLM-sized batches, never splitting a block."""
    batches: list[list[Block]] = []
    current: list[Block] = []
    size = 0
    for loc, text in blocks:
        if current and size + len(text) > max_chars:
            batches.append(current)
            current, size = [], 0
        current.append((loc, text))
        size += len(text)
    if current:
        batches.append(current)
    return batches


def render(blocks: list[Block]) -> str:
    """Format blocks for a prompt with their locators attached."""
    return "\n\n".join(f"[[{loc}]]\n{text}" for loc, text in blocks)
